import streamlit as st
import fitz  # PyMuPDF
from pptx import Presentation
from pptx.util import Inches
from PIL import Image
import tempfile
import os

st.set_page_config(page_title="PDF → PowerPoint Converter", layout="centered")

st.title("📑➡️📊 PDF to PowerPoint Converter")
st.write("Upload one or more PDFs. Each PDF page will become a slide with the **same proportions as the PDF**.")

uploaded_files = st.file_uploader(
    "Choose PDF files",
    type=["pdf"],
    accept_multiple_files=True
)

dpi = st.slider("Image DPI (higher = better quality, bigger file)", 100, 400, 200)

def pdf_to_pptx(pdf_path, dpi=200):
    # Open PDF with PyMuPDF
    doc = fitz.open(pdf_path)
    first_page = doc[0]
    rect = first_page.rect

    # Convert PDF page size from points (1/72 in) to inches
    pdf_width_in = rect.width / 72
    pdf_height_in = rect.height / 72

    # Create PowerPoint with same proportions
    prs = Presentation()
    prs.slide_width = Inches(pdf_width_in)
    prs.slide_height = Inches(pdf_height_in)

    for page in doc:
        # Render page to image at desired DPI
        mat = fitz.Matrix(dpi / 72, dpi / 72)
        pix = page.get_pixmap(matrix=mat)
        img = Image.frombytes("RGB", [pix.width, pix.height], pix.samples)

        # Save temp PNG for inserting into slide
        with tempfile.NamedTemporaryFile(delete=False, suffix=".png") as tmp_img:
            img.save(tmp_img.name, "PNG")
            slide = prs.slides.add_slide(prs.slide_layouts[6])
            slide.shapes.add_picture(tmp_img.name, 0, 0, width=prs.slide_width, height=prs.slide_height)
            os.remove(tmp_img.name)

    # Save PPTX to temp file
    out_pptx = tempfile.NamedTemporaryFile(delete=False, suffix=".pptx")
    prs.save(out_pptx.name)
    return out_pptx.name

if uploaded_files:
    for uploaded_file in uploaded_files:
        with st.spinner(f"Processing {uploaded_file.name}..."):
            with tempfile.NamedTemporaryFile(delete=False, suffix=".pdf") as tmp_pdf:
                tmp_pdf.write(uploaded_file.read())
                tmp_pdf_path = tmp_pdf.name

            pptx_path = pdf_to_pptx(tmp_pdf_path, dpi=dpi)

            st.success(f"Finished: {uploaded_file.name}")
            st.download_button(
                label=f"⬇️ Download {uploaded_file.name.replace('.pdf','.pptx')}",
                data=open(pptx_path, "rb").read(),
                file_name=uploaded_file.name.replace(".pdf", ".pptx"),
                mime="application/vnd.openxmlformats-officedocument.presentationml.presentation"
            )

            os.remove(tmp_pdf_path)
            os.remove(pptx_path)
